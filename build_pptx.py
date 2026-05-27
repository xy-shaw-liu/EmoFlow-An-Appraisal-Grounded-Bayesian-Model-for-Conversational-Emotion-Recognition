"""Generate EmoFlow_Presentation.pptx — ECS 271 final presentation deck.

Run: python build_pptx.py  →  EmoFlow_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW, SH = prs.slide_width, prs.slide_height

# Palette
NAVY = RGBColor(0x0B, 0x2B, 0x4A)
ACCENT = RGBColor(0xE2, 0x6D, 0x5C)
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
GREEN = RGBColor(0x2F, 0x9E, 0x44)
RED = RGBColor(0xC9, 0x2A, 0x2A)

BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, page_num, total=24):
    # Top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    # Title
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.7),
             size=28, bold=True, color=NAVY)
    # Page indicator
    add_text(slide, f"{page_num} / {total}",
             Inches(12.0), Inches(0.25), Inches(1.0), Inches(0.4),
             size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    # Divider
    add_rect(slide, Inches(0.5), Inches(1.0), Inches(12.33), Emu(12000), MUTED)


def add_footer(slide):
    add_text(slide, "EmoFlow  ·  ECS 271 Spring 2026",
             Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             size=10, color=MUTED)


def bullets(slide, items, left, top, width, height, size=18, gap=0.08):
    """items: list[str | (str, level)]. level 0=main, 1=sub."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, level = it
        else:
            text, level = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(size * gap * 10)
        if level == 0:
            bullet = "•  "
            indent = 0
            sz = size
            col = INK
            b = False
        else:
            bullet = "–  "
            indent = Inches(0.35)
            sz = size - 2
            col = MUTED
            b = False
        p.level = level
        r = p.add_run()
        r.text = bullet + text
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = col
        r.font.name = "Calibri"
    return box


def add_table(slide, data, left, top, width, height, header=True,
              col_widths=None, font_size=14, header_fill=NAVY, header_color=RGBColor(0xFF, 0xFF, 0xFF),
              highlight_rows=None, highlight_color=RGBColor(0xFF, 0xF4, 0xE6)):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = int(width * w / total)
    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = tf.paragraphs[0].add_run()
            r.text = str(data[ri][ci])
            r.font.size = Pt(font_size)
            r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                r.font.bold = True
                r.font.color.rgb = header_color
            elif highlight_rows and ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_color
                r.font.bold = True
                r.font.color.rgb = INK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.color.rgb = INK
    return tbl


# ---------------------------------------------------------------------------
# Slide 0: Title (not counted in 24)
# ---------------------------------------------------------------------------
s = add_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.0), SW, Inches(0.08), ACCENT)
add_text(s, "EmoFlow",
         Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.4),
         size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, "An Appraisal-Grounded Bayesian Model\nfor Conversational Emotion Recognition",
         Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.4),
         size=28, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "Xiaoyi Liu   ·   Zhiye Jiang   ·   Ruitian Liu",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
         size=22, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "ECS 271 — Machine Learning — UC Davis, Spring 2026",
         Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
         size=16, color=RGBColor(0xC0, 0xCB, 0xDC), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Section transition helper (also not counted)
# ---------------------------------------------------------------------------
def section_slide(num, title, subtitle=""):
    s = add_slide()
    add_rect(s, 0, 0, SW, SH, LIGHT)
    add_rect(s, Inches(0.5), Inches(2.8), Inches(0.15), Inches(1.5), ACCENT)
    add_text(s, f"PART {num}", Inches(0.8), Inches(2.8), Inches(6), Inches(0.6),
             size=18, bold=True, color=ACCENT)
    add_text(s, title, Inches(0.8), Inches(3.3), Inches(12), Inches(1.2),
             size=44, bold=True, color=NAVY)
    if subtitle:
        add_text(s, subtitle, Inches(0.8), Inches(4.6), Inches(12), Inches(0.6),
                 size=18, color=MUTED)
    return s


# ===========================================================================
# PART I — Background & Problem
# ===========================================================================
section_slide("I", "Background & Problem", "What is conversational emotion recognition, and why is it hard?")

# ---- Slide 1: The Task --------------------------------------------------
s = add_slide(); add_header(s, "The Task: Conversational Emotion Recognition (ERC)", 1)
add_text(s, "Given a dialogue, predict the emotion of every turn.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=20, color=INK, bold=True)

# Example dialogue box
box = add_rect(s, Inches(0.5), Inches(1.9), Inches(7.5), Inches(4.5), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s, "Example — same utterance, three emotions",
         Inches(0.7), Inches(2.0), Inches(7), Inches(0.4),
         size=14, bold=True, color=MUTED)
add_text(s,
         "A:  \"Did you finish the report?\"\n"
         "B:  \"Sure.\"                           → neutral  (acknowledgement)\n\n"
         "A:  \"You owe me twenty bucks.\"\n"
         "B:  \"Sure.\"                           → disgust  (sarcastic)\n\n"
         "A:  \"…so will you marry me?\"\n"
         "B:  \"Sure!\"                           → surprise / joy",
         Inches(0.7), Inches(2.5), Inches(7.2), Inches(3.7),
         size=15, color=INK, font="Consolas")

# Right side — formal definition
add_text(s, "Formal definition",
         Inches(8.3), Inches(1.9), Inches(4.5), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Input: utterances u₁, u₂, …, u_T from one or more speakers",
    "Output: emotion label y_t ∈ ℰ for every turn",
    "ℰ = {neutral, joy, sadness, anger, fear, disgust, surprise}",
    "Benchmark: MELD (Friends transcripts, 13K turns, 7 classes)",
], Inches(8.3), Inches(2.4), Inches(4.7), Inches(4.0), size=14)
add_footer(s)

# ---- Slide 2: Three structural challenges -------------------------------
s = add_slide(); add_header(s, "Why ERC Is Harder Than Sentence-Level Emotion", 2)
challenge_x = [Inches(0.5), Inches(4.85), Inches(9.2)]
labels = [
    ("Context dependence", "The same words mean different\nthings in different histories.",
     "\"Sure.\" depends entirely on the\nturns that came before."),
    ("Emotional inertia", "Emotions persist across turns.",
     "After three angry turns, a flat\n\"Whatever.\" is still angry."),
    ("Severe class imbalance", "Long-tailed label distribution.",
     "MELD: neutral 48 %, fear + disgust\ntogether < 5 %."),
]
for i, (h, sub, ex) in enumerate(labels):
    add_rect(s, challenge_x[i], Inches(1.5), Inches(3.8), Inches(5), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
    add_rect(s, challenge_x[i], Inches(1.5), Inches(3.8), Inches(0.08), ACCENT)
    add_text(s, f"({i+1})", challenge_x[i] + Inches(0.3), Inches(1.7), Inches(1), Inches(0.5),
             size=14, bold=True, color=ACCENT)
    add_text(s, h, challenge_x[i] + Inches(0.3), Inches(2.1), Inches(3.4), Inches(0.6),
             size=20, bold=True, color=NAVY)
    add_text(s, sub, challenge_x[i] + Inches(0.3), Inches(2.9), Inches(3.4), Inches(1.4),
             size=15, color=INK)
    add_text(s, "Example", challenge_x[i] + Inches(0.3), Inches(4.4), Inches(3.4), Inches(0.4),
             size=12, bold=True, color=MUTED)
    add_text(s, ex, challenge_x[i] + Inches(0.3), Inches(4.7), Inches(3.4), Inches(1.7),
             size=13, color=INK, font="Consolas")
add_text(s, "Naive classifiers collapse to the majority class.   Generic recurrent models overfit small ERC datasets.",
         Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s)

# ---- Slide 3: Limitations of existing approaches ------------------------
s = add_slide(); add_header(s, "Limitations of Existing ERC Models", 3)
add_text(s, "State-of-the-art models work, but they are black boxes.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=18, bold=True, color=INK)

data = [
    ["Family", "Example", "Issue we want to address"],
    ["RNN-based", "DialogueRNN (2019)", "Context, speaker, emotion entangled in one hidden state"],
    ["Graph-based", "DialogueGCN, MMGCN", "Heavy parameters; no interpretable intermediate"],
    ["LLM fine-tune", "COSMIC (BERT-large)", "100M – 300M+ parameters; opaque representations"],
    ["Pure prompting", "Zero-/few-shot LLM", "No temporal model; unstable on rare classes"],
]
add_table(s, data, Inches(0.5), Inches(1.9), Inches(12.3), Inches(2.7),
          col_widths=[1.4, 2.2, 5.5], font_size=15)

add_text(s, "What's missing",
         Inches(0.5), Inches(4.9), Inches(12), Inches(0.4),
         size=18, bold=True, color=NAVY)
bullets(s, [
    "An interpretable bottleneck that humans can inspect",
    "An explicit, structured way to combine context with current evidence",
    "Strong inductive bias to survive the long-tailed label distribution",
], Inches(0.5), Inches(5.4), Inches(12.3), Inches(2.0), size=17)
add_footer(s)

# ---- Slide 4: Our angle -------------------------------------------------
s = add_slide(); add_header(s, "Our Angle: Appraisal Theory + Bayesian Fusion", 4)
add_text(s, "Borrow structure from emotion psychology rather than scaling parameters.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=18, bold=True, color=INK)

# Three-component preview
boxes = [
    ("1.  Appraisal grounding",
     "Force the bottleneck to be an 8-dim\nScherer appraisal vector — each axis\nhas psychological meaning.",
     "→ Interpretability + theoretical prior"),
    ("2.  Bayesian fusion",
     "Keep context prior and current\nlikelihood separable instead of\nentangled in an RNN state.",
     "→ Inductive bias + ablatable structure"),
    ("3.  Exponential-decay memory",
     "Aggregate past appraisals with a\nsingle learnable λ rather than a\ngeneric recurrent cell.",
     "→ Resists overfitting on small data"),
]
for i, (h, body, take) in enumerate(boxes):
    x = Inches(0.5 + i * 4.35)
    add_rect(s, x, Inches(1.9), Inches(4.1), Inches(4.7), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
    add_rect(s, x, Inches(1.9), Inches(4.1), Inches(0.6), NAVY)
    add_text(s, h, x + Inches(0.2), Inches(2.0), Inches(3.7), Inches(0.5),
             size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, body, x + Inches(0.2), Inches(2.7), Inches(3.7), Inches(2.5),
             size=14, color=INK)
    add_text(s, take, x + Inches(0.2), Inches(5.7), Inches(3.7), Inches(0.8),
             size=13, bold=True, color=ACCENT)

add_text(s, "Headline:  wF1 = 0.62 on MELD with only ~20 M trainable params — controlled ablations (Part III) isolate the architecture's gains.",
         Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_footer(s)


# ===========================================================================
# PART II — Method
# ===========================================================================
section_slide("II", "Method", "Three components, grounded in Scherer's Component Process Model.")

# ---- Slide 5: Scherer's CPM + 8 dims -----------------------------------
s = add_slide(); add_header(s, "Theoretical Foundation: Scherer's Component Process Model", 5)
add_text(s, "Emotions arise from a sequence of cognitive appraisal checks of an event.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

# 8 appraisal dimensions
dims = [
    ("expectedness", "Was the event expected?"),
    ("unpleasantness", "Is it intrinsically unpleasant?"),
    ("goal_hindrance", "Does it block my goals?"),
    ("external_causation", "Caused by someone / something else?"),
    ("coping_potential", "Can I cope with / control it?"),
    ("unfairness", "Is it unfair?"),
    ("immorality", "Does it violate moral norms?"),
    ("self_consistency", "Does it fit my self-image?"),
]
add_text(s, "The 8 ISEAR appraisal dimensions",
         Inches(0.5), Inches(1.85), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
for i, (d, q) in enumerate(dims):
    col, row = i % 2, i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.35 + row * 0.6)
    add_text(s, f"  {d}", x, y, Inches(2.6), Inches(0.4),
             size=14, bold=True, color=ACCENT, font="Consolas")
    add_text(s, q, x + Inches(2.6), y, Inches(3.6), Inches(0.4),
             size=13, color=INK)
add_text(s,
         "Each emotion has a distinctive appraisal signature.\n"
         "Example — anger:  max-valued on unpleasantness, goal_hindrance, unfairness, immorality.",
         Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0),
         size=15, color=INK)
add_text(s, "Reference:  Scherer, K. R. (2001). Appraisal Processes in Emotion — Table 5.5.",
         Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
         size=12, color=MUTED)
add_footer(s)

# ---- Slide 6: Architecture overview ------------------------------------
s = add_slide(); add_header(s, "Architecture Overview", 6)
add_text(s, "Three stages mirroring CPM:  stimulus → appraisal → emotion.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

# Build pipeline diagram
def stage_box(left, top, w, h, title, sub, fill=NAVY):
    add_rect(s, left, top, w, h, fill)
    add_text(s, title, left, top + Inches(0.15), w, Inches(0.4),
             size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, sub, left, top + Inches(0.55), w, h - Inches(0.6),
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)

def arrow(left, top, w, label=""):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, Inches(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = ACCENT; a.line.fill.background()
    if label:
        add_text(s, label, left, top - Inches(0.4), w, Inches(0.35),
                 size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")

stage_box(Inches(0.4), Inches(2.6), Inches(2.7), Inches(1.6),
          "StimulusEncoder", "frozen LLaMA-3-8B\n+ LoRA + AppraisalHead\nencoder.py")
arrow(Inches(3.2), Inches(3.0), Inches(0.9), "appraisal\n(8-d)")
stage_box(Inches(4.2), Inches(2.6), Inches(2.7), Inches(1.6),
          "TemporalMemory", "exp(−λ·Δt) decay\nover past appraisals\nmemory.py")
arrow(Inches(7.0), Inches(3.0), Inches(0.9), "memory\nstate")
stage_box(Inches(8.0), Inches(2.6), Inches(2.7), Inches(1.6),
          "BayesianHead", "log-additive fusion\nprior + likelihood\nbayes.py")
arrow(Inches(10.8), Inches(3.0), Inches(0.9), "6 logits")
stage_box(Inches(11.8), Inches(2.6), Inches(1.3), Inches(1.6),
          "Output", "sigmoid + τ\n→ 6 emotions\nor neutral", fill=ACCENT)

# Side branch — likelihood path
add_rect(s, Inches(0.4), Inches(4.6), Inches(10.4), Emu(8000), MUTED)
add_text(s,
         "Likelihood path:  current appraisal → likelihood logits  (BayesianHead also takes the raw current appraisal directly)",
         Inches(0.4), Inches(4.7), Inches(11), Inches(0.4),
         size=12, color=MUTED, align=PP_ALIGN.LEFT, font="Consolas")

# Trainable budget
add_text(s, "Trainable parameter budget", Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "AppraisalHead — Linear(4096→4096→8):  ~17 M  (the bulk)",
    "LoRA adapters — r = 8 on q_proj / v_proj × 32 layers:  ~4 M",
    "TemporalMemory (1 scalar λ) + BayesianHead (two 8→64→6 MLPs):  ~2 K",
    "Total trainable:  ~20 M     ·     0.25% of the frozen 8 B backbone",
], Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.5), size=14)
add_footer(s)

# ---- Slide 7: StimulusEncoder ------------------------------------------
s = add_slide(); add_header(s, "Component 1 — StimulusEncoder", 7)
add_text(s, "Map an utterance to an 8-dim appraisal vector.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

add_text(s, "Pipeline",
         Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Tokenize utterance with LLaMA-3 tokenizer",
    "Forward through frozen LLaMA-3-8B + LoRA adapters",
    "Pool the last non-pad token's hidden state (∈ ℝ⁴⁰⁹⁶)",
    "AppraisalHead: Linear(4096→4096) → GELU → Linear(4096→8)",
    "Output: 8-d vector a_t — one entry per Scherer dimension",
], Inches(0.5), Inches(2.15), Inches(6.5), Inches(3.5), size=14)

add_text(s, "Design choices",
         Inches(7.2), Inches(1.7), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "LoRA only (r=8, α=16) — backbone weights frozen",
    "QLoRA (NF4 4-bit) backbone — fits in a 24 GB GPU",
    "DistilBERT swap for fast local dev / ablation",
    "No sigmoid on the head output  (see Part III)",
    "Weak supervision: MSE against Scherer prototypes + emotion BCE downstream",
], Inches(7.2), Inches(2.15), Inches(5.7), Inches(4.0), size=14)

add_text(s, "Trainable here:  AppraisalHead (~17 M)  +  LoRA adapters (~4 M)  =  the ~20 M budget.",
         Inches(0.5), Inches(6.65), Inches(12), Inches(0.4),
         size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s)

# ---- Slide 8: Appraisal targets ----------------------------------------
s = add_slide(); add_header(s, "Appraisal Targets — Raw Z-scores → [0, 1]", 8)
add_text(s, "Final per-emotion prototype vectors (min-max normalized) used as weak-supervision MSE targets.",
         Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.45),
         size=15, color=INK)

# Final normalized targets (from appraisal_targets.py, 2 d.p.)
data = [
    ["Emotion",  "expect", "unplea", "goal_h", "ext_c", "cope", "unfair", "immor", "self_c"],
    ["joy",      "1.00",   "0.00",   "0.00",   "0.00",  "1.00", "0.00",   "0.00",  "1.00"],
    ["fear",     "0.54",   "0.97",   "0.84",   "0.60",  "0.08", "0.55",   "0.63",  "0.08"],
    ["anger",    "0.48",   "1.00",   "1.00",   "0.42",  "0.53", "1.00",   "1.00",  "0.06"],
    ["sadness",  "0.63",   "0.98",   "0.96",   "1.00",  "0.00", "0.64",   "0.55",  "0.00"],
    ["disgust",  "0.49",   "1.00",   "0.85",   "0.67",  "0.47", "0.76",   "1.00",  "0.11"],
    ["surprise", "0.00",   "0.83",   "0.76",   "0.22",  "0.44", "0.55",   "0.69",  "0.13"],
]
add_table(s, data, Inches(0.5), Inches(1.65), Inches(12.3), Inches(2.45),
          col_widths=[1.6, 1.0, 1.0, 1.0, 0.95, 0.9, 1.0, 0.95, 1.0], font_size=12)

# Preprocessing band
add_text(s, "Preprocessing — why the raw table can't be used as-is",
         Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
         size=16, bold=True, color=NAVY)

add_text(s, "Raw data (Scherer Table 5.5)",
         Inches(0.5), Inches(4.75), Inches(6), Inches(0.35),
         size=14, bold=True, color=ACCENT)
bullets(s, [
    "Empirical Z-scores from a 37-country ISEAR study (n ≈ 3000).",
    "Centered on the population mean = 0  →  values are signed.",
    "e.g.  joy·unpleasantness = −2.00,   anger·unfairness = +0.58.",
], Inches(0.5), Inches(5.1), Inches(6.1), Inches(2.0), size=13)

add_text(s, "Transform → [0, 1]",
         Inches(6.9), Inches(4.75), Inches(6), Inches(0.35),
         size=14, bold=True, color=ACCENT)
bullets(s, [
    "Per-dimension min-max across the 6 emotions.",
    "Per-dim (not global) equalizes scales so no high-variance dim dominates the MSE loss — while preserving emotion ordering within each dim.",
    "Surprise (absent from ISEAR) derived: expectedness set below the empirical min → 0 = max novelty; other dims = mean.",
], Inches(6.9), Inches(5.1), Inches(6.0), Inches(2.0), size=13)
add_footer(s)

# ---- Slide 9: TemporalMemory -------------------------------------------
s = add_slide(); add_header(s, "Component 2 — TemporalMemory", 9)
add_text(s, "Aggregate past appraisals with a single learnable decay constant.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

# Equation
add_rect(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.7), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s, "Memory update  (normalized weighted average)",
         Inches(0.7), Inches(1.95), Inches(8), Inches(0.4),
         size=14, bold=True, color=MUTED)
add_text(s,
         "h_t  =  ( Σ_{i ≤ t}  w_i · a_i )  /  ( Σ_{i ≤ t} w_i ) ,      w_i = exp(−λ · (t − i))",
         Inches(0.7), Inches(2.35), Inches(12), Inches(0.5),
         size=17, color=INK, font="Consolas")
add_text(s,
         "Implemented as softmax(−λ·Δt) with causal masking.   λ = softplus(θ), learnable, init 0.1.\n"
         "λ = 0  →  uniform mean over past + current;    λ → ∞  →  only the current turn.",
         Inches(0.7), Inches(2.9), Inches(12), Inches(0.6),
         size=13, color=MUTED, font="Consolas")

add_text(s, "Why exponential decay (not BiLSTM)",
         Inches(0.5), Inches(3.75), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "One learnable scalar  vs.  thousands of recurrent weights → strong inductive bias.",
    "Captures the right qualitative shape: recent turns matter most, distant turns fade.",
    "Closed-form, parallelizable, no vanishing-gradient issues.",
    "Ablation (slide 20):  BiLSTM memory is 14 wF1 points WORSE than no memory at all.",
], Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), size=15)
add_footer(s)

# ---- Slide 10: BayesianHead --------------------------------------------
s = add_slide(); add_header(s, "Component 3 — BayesianHead", 10)
add_text(s, "Fuse context-driven prior with current-evidence likelihood — log-additively.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

add_rect(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(2.2), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s,
         "log P(y_t | a_t, h_t)   ∝   log P(y_t | h_t)         +   log P(a_t | y_t)\n"
         "                                  prior  (from memory)         likelihood  (current)",
         Inches(0.7), Inches(2.05), Inches(12), Inches(1.2),
         size=18, color=INK, font="Consolas")
add_text(s,
         "Implementation:   logits = prior_head(h_t)  +  likelihood_head(a_t)\n"
         "   each head:  Linear(8→64) → GELU → Dropout(0.1) → Linear(64→6)   (K = 6 non-neutral classes)",
         Inches(0.7), Inches(3.05), Inches(12), Inches(0.9),
         size=13, color=MUTED, font="Consolas")

add_text(s, "Why two heads instead of one big head",
         Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Decoupling forces interpretable structure — each MLP has an isolated job.",
    "Ablation-friendly: replace MLP_prior with a uniform vector → measures context contribution directly.",
    "Mathematically equivalent to an exact Dirichlet–Categorical posterior under uniform marginals.",
    "Adds structure almost for free: each branch is a tiny 8→64→6 MLP (~1 K weights).",
], Inches(0.5), Inches(4.75), Inches(12.3), Inches(2.3), size=14)
add_footer(s)

# ---- Slide 11: Joint loss + multilabel reformulation -------------------
s = add_slide(); add_header(s, "Training Objective + Multilabel Reformulation", 11)

add_text(s, "Joint loss",
         Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.3), Inches(1.6), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s,
         "L  =  L_emo  +  α · L_app",
         Inches(0.7), Inches(1.85), Inches(6), Inches(0.5),
         size=18, color=INK, bold=True, font="Consolas")
add_text(s,
         "L_emo:  BCE-with-logits over 7 emotions   (multilabel)\n"
         "L_app:   MSE  â_t  vs  Scherer prototype[y_t]\n"
         "α = 0.1   (tuned on dev)",
         Inches(0.7), Inches(2.3), Inches(6), Inches(0.9),
         size=12, color=MUTED, font="Consolas")

add_text(s, "Multilabel reformulation",
         Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Treat neutral as the all-zero target — not a class.",
    "Every non-neutral emotion is an independent 0/1 head.",
    "Structurally prevents the model from collapsing to \"all neutral\".",
    "Pivotal:  broke us out of the all-neutral failure mode (§III).",
], Inches(7.0), Inches(1.6), Inches(6.0), Inches(3.3), size=14)

# Bottom
add_text(s, "Optimization & decoding",
         Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "AdamW,  lr = 5e-4,  batch size 2 dialogues,  3 epochs.   bf16 backbone / fp32 heads.",
    "Dialogue-level WeightedRandomSampler for rare-class oversampling.",
    "Single threshold τ swept on dev over {0.2 … 0.6};  τ = 0.2 selected — predict neutral if max sigmoid prob < τ.",
], Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.8), size=14)
add_footer(s)

# ---- Slide 12: Scope pivot ---------------------------------------------
s = add_slide(); add_header(s, "Scope: Proposal vs. Delivered", 12)
add_text(s, "We pivoted from a 6-module generation pipeline to a focused, finished classifier.",
         Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5),
         size=17, color=INK)

data = [
    ["Proposal (Feb 2026)", "Delivered (May 2026)"],
    ["6 modules: …→ classifier → Response Generator", "Classifier only — generator deferred to future work"],
    ["Primary dataset: EmoryNLP", "Primary dataset: MELD  (EmoryNLP class-collapsed pre-fix)"],
    ["5 aggregated appraisal dimensions", "8 raw ISEAR dimensions (preserves anger / disgust separation)"],
    ["Exact Dirichlet–Categorical update", "Log-additive MLP fusion  (equivalent under uniform marginals)"],
    ["LoRA r = 16", "LoRA r = 8  (QLoRA memory constraint)"],
    ["Human Coherence Score evaluation", "Not performed — tied to absent Response Generator"],
]
add_table(s, data, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.2),
          col_widths=[5.5, 6.8], font_size=14)

add_text(s,
         "Reason for the pivot:  failure-mode debugging (next section) consumed the compute budget originally allocated to the generator.\n"
         "The classifier is a strict prerequisite for the generator — and a self-contained contribution.",
         Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0),
         size=14, color=MUTED, align=PP_ALIGN.LEFT)
add_footer(s)


# ===========================================================================
# PART III — Experiments & Analysis
# ===========================================================================
section_slide("III", "Experiments & Analysis", "Data, a failure story, results, and ablations.")

# ---- Slide 13: Data -----------------------------------------------------
s = add_slide(); add_header(s, "Data — MELD + Selective Cross-Dataset Augmentation", 13)

# Datasets
add_text(s, "Datasets",
         Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
data = [
    ["Dataset", "Role", "Train dlg"],
    ["MELD",         "Primary benchmark",         "1,038"],
    ["DailyDialog",  "Rare-class augmentation",   "11,118"],
    ["EmoryNLP",     "Secondary (limited use)",   "713"],
]
add_table(s, data, Inches(0.5), Inches(1.6), Inches(6.2), Inches(2.0),
          col_widths=[2.0, 3.2, 1.2], font_size=14)

# Class imbalance
add_text(s, "MELD train distribution",
         Inches(7.0), Inches(1.15), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
classes = [("neutral", 47), ("joy", 17), ("surprise", 12), ("anger", 11),
           ("sadness", 7), ("disgust", 3), ("fear", 3)]
bar_x = Inches(7.0)
bar_top = Inches(1.65)
max_w = Inches(5.5)
for i, (lab, pct) in enumerate(classes):
    y = bar_top + Inches(i * 0.35)
    add_text(s, f"{lab}", bar_x, y, Inches(1.2), Inches(0.3), size=12, color=INK)
    w = int(max_w * pct / 50)
    color = ACCENT if pct < 10 else NAVY
    add_rect(s, bar_x + Inches(1.2), y + Inches(0.05), w, Inches(0.2), color)
    add_text(s, f"{pct} %", bar_x + Inches(1.2) + w + Inches(0.05), y, Inches(1.0), Inches(0.3),
             size=12, color=MUTED)

add_text(s, "Augmentation strategy — DD-rare3",
         Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Add only DailyDialog dialogues containing a fear / disgust / sadness turn (the 3 rarest).",
    "Naive full merge backfires: DailyDialog is 83% no_emotion → would push neutral 47% → 66%.",
    "Selective merge raises rare-class share 12.2% → 15.0%;  fear train count 268 → 414.",
    "Then dialogue-level WeightedRandomSampler oversamples rare-class dialogues up to 30×.",
], Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4), size=14)
add_footer(s)

# ---- Slide 14: Six failed mitigations ----------------------------------
s = add_slide(); add_header(s, "Failure Analysis — Six Mitigations, One Hidden Cause", 14)
add_text(s, "Six imbalance-mitigation attempts — none recovered the model.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

data = [
    ["#", "Intervention",                       "Result",     "Failure mode"],
    ["1", "Inverse-freq class weights",         "wF1 0.31",   "predict all neutral"],
    ["2", "+ Label smoothing (ε = 0.1)",        "wF1 0.31",   "predict all neutral"],
    ["3", "+ WeightedRandomSampler",            "0.25 dev",   "predict all neutral"],
    ["4", "Multilabel BCE (5-d appraisal)",     "0.25 dev",   "marginal collapse"],
    ["5", "+ Per-dim pos_weight (1−p)/p",       "0.0025 dev", "oversteer to rare"],
    ["6", "5 → 8-dim appraisal expansion",      "0.25 dev",   "collapse persists"],
]
add_table(s, data, Inches(0.5), Inches(1.85), Inches(9.0), Inches(3.6),
          col_widths=[0.5, 4.3, 1.6, 2.6], font_size=13,
          highlight_rows=list(range(1, 7)))

add_text(s, "🚩  Attempts 1–3 bit-identical", Inches(9.8), Inches(1.85), Inches(3.4), Inches(0.5),
         size=15, bold=True, color=RED)
bullets(s, [
    "Dev wF1 = 0.2523, 194/607 transitions correct — identical across all three.",
    "Same numbers regardless of the loss-side fix.",
    "First hint: the bug is upstream of the loss.",
], Inches(9.8), Inches(2.4), Inches(3.4), Inches(3.0), size=12)

add_text(s, "Lesson:  identical metrics across different loss-side interventions ⇒ the representation, not the loss, is broken.",
         Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_footer(s)

# ---- Slide 15: Root cause — sigmoid saturation -------------------------
s = add_slide(); add_header(s, "Root Cause — Sigmoid Saturation in the Appraisal Head", 15)
add_text(s, "We diagnosed it by checking representations, not metrics.",
         Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
         size=17, color=INK)

add_text(s, "Diagnostic",
         Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Forward 3 semantically distinct sentences through the encoder.",
    "Inspect each appraisal dim's std across the 3 inputs.",
    "Output bit-identical:  [1,1,1,0,0,1,1,1];  std = 0.000 on every dim.",
    "Zero discrimination — the bottleneck's effective input dimension is 0.",
], Inches(0.5), Inches(2.15), Inches(6.0), Inches(2.6), size=14)

add_text(s, "Why",
         Inches(7.0), Inches(1.7), Inches(6), Inches(0.4),
         size=16, bold=True, color=NAVY)
add_rect(s, Inches(7.0), Inches(2.15), Inches(6.0), Inches(2.6), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s,
         "AppraisalHead ended with  σ(·)  to match\n"
         "Scherer prototypes in [0, 1].\n\n"
         "Pre-activation values drifted to large |x|.\n"
         "→  σ saturates  →  σ' ≈ 0\n"
         "→  zero gradient through the head\n"
         "→  encoder learns nothing.",
         Inches(7.2), Inches(2.3), Inches(5.8), Inches(2.3),
         size=14, color=INK)

add_text(s, "Why every imbalance fix failed",
         Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Class weights, focal loss, oversampling all modify the emotion-CE gradient.",
    "But that gradient flows back through a saturated sigmoid into the encoder.",
    "The encoder cannot update → model stays at the prior (majority class).",
    "All six interventions produced bit-identical metrics for this reason.",
], Inches(0.5), Inches(5.45), Inches(12.3), Inches(2.0), size=13)
add_footer(s)

# ---- Slide 16: The fix --------------------------------------------------
s = add_slide(); add_header(s, "The Fix — One Line", 16)

add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.5), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s, "Before",
         Inches(0.8), Inches(1.55), Inches(5), Inches(0.4),
         size=14, bold=True, color=RED)
add_text(s,
         "a_t = torch.sigmoid(self.head(h))   # forces output to [0, 1]",
         Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6),
         size=18, color=INK, font="Consolas")
add_text(s, "After",
         Inches(0.8), Inches(2.7), Inches(5), Inches(0.4),
         size=14, bold=True, color=GREEN)
add_text(s,
         "a_t = self.head(h)                  # unbounded; range enforced by MSE loss",
         Inches(0.8), Inches(3.15), Inches(11.5), Inches(0.6),
         size=18, color=INK, font="Consolas")

add_text(s, "Effect",
         Inches(0.5), Inches(4.15), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Encoder output std across inputs:  0.000  →  0.413.",
    "Post-fix appraisals are input-dependent and Scherer-aligned (e.g. surprise → lowest expectedness).",
    "Imbalance mitigations now actually reach the encoder — previously blocked by the saturated head.",
    "wF1:  ~0.25 (collapsed)  →  0.42 (post-fix)  →  0.62 (+ DD-rare3 augmentation).",
], Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.0), size=15)

add_text(s,
         "Generalizable lesson:  prefer unbounded activations at internal low-dim bottlenecks; "
         "express target ranges through the loss, not the architecture.",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s)

# ---- Slide 17: Main results --------------------------------------------
s = add_slide(); add_header(s, "Main Results — Ours vs. Our Baselines", 17)
add_text(s, "Same config (multilabel BCE, α = 0.1, MELD + DD-rare3, oversample, no sigmoid) — only the model class differs.",
         Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5),
         size=15, color=INK)

data = [
    ["Model",                          "wF1",      "wF1 (6-way)", "ETA",     "trainable"],
    ["LSTM (BiLSTM over appraisals)",  "0.4241",   "0.2719",      "0.4095",  "~20 M"],
    ["Stateless (no memory)",          "0.5631",   "0.4553",      "0.5462",  "~20 M"],
    ["EmoFlow (λ = 0, no decay)",      "0.6052",   "0.5348",      "0.5617",  "~20 M"],
    ["EmoFlow (learned λ)",            "0.6171",   "0.5312",      "0.5528",  "~20 M"],
]
add_table(s, data, Inches(0.5), Inches(1.95), Inches(12.3), Inches(3.0),
          col_widths=[4.5, 1.7, 2.1, 1.5, 2.0], font_size=15,
          highlight_rows=[4])

add_text(s, "Takeaways",
         Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "EmoFlow (learned λ):  +19 wF1 over the LSTM baseline at the same parameter budget.",
    "ETA (Emotion Transition Accuracy): tests turn-to-turn changes — our hardest sub-metric.",
    "wF1_6way drops the dominant neutral class — shows we are actually learning rare classes.",
], Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.6), size=14)
add_footer(s)

# ---- Slide 18: Comparison with published models ------------------------
s = add_slide(); add_header(s, "Comparison with Published Models on MELD", 18)

data = [
    ["Model",          "Year", "MELD wF1", "Setting / trainable params"],
    ["DialogueRNN",    "2019", "0.57",     "text · ~3 M"],
    ["DialogueGCN",    "2019", "0.58",     "text · ~5 M"],
    ["MMGCN †",        "2021", "0.59",     "multimodal · ~110 M+"],
    ["EmoFlow (ours)", "2026", "0.62",     "text · ~20 M trainable"],
    ["COSMIC",         "2020", "0.65",     "text · BERT-large ~340 M"],
]
add_table(s, data, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.4),
          col_widths=[3.3, 1.2, 1.6, 6.2], font_size=16,
          highlight_rows=[4])

add_text(s, "How to read this  —  indicative, NOT controlled",
         Inches(0.5), Inches(5.15), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Raw MELD number: EmoFlow 0.62 sits above DialogueRNN / GCN / MMGCN, below COSMIC.",
    "BUT EmoFlow adds DailyDialog rare-class augmentation; baselines train on MELD only.",
    "MELD-only EmoFlow ≈ 0.42 — the augmentation drives most of the gap, not the architecture alone.",
    "Controlled evidence = the equal-data, equal-param ablation (previous slide).",
    "† MMGCN's quoted 0.66 is IEMOCAP; its MELD score is 0.5865 · DialogueRNN MELD = 0.5703.",
], Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.9), size=12)
add_footer(s)

# ---- Slide 19: Per-class + fear ----------------------------------------
s = add_slide(); add_header(s, "Per-Class Breakdown — and the Fear Problem", 19)

data = [
    ["Emotion",  "Support", "Precision", "Recall", "F1"],
    ["neutral",  "1256",    "0.75",      "0.78",   "0.76"],
    ["joy",      "402",     "0.62",      "0.54",   "0.58"],
    ["surprise", "281",     "0.50",      "0.65",   "0.57"],
    ["anger",    "345",     "0.53",      "0.48",   "0.51"],
    ["sadness",  "208",     "0.51",      "0.24",   "0.32"],
    ["disgust",  "68",      "0.17",      "0.44",   "0.25"],
    ["fear",     "50",      "0.00",      "0.00",   "0.00"],
]
add_table(s, data, Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.8),
          col_widths=[1.8, 1.3, 1.5, 1.3, 1.2], font_size=14,
          highlight_rows=[7])

add_text(s, "Six of seven classes learned non-trivially.",
         Inches(8.3), Inches(1.5), Inches(4.8), Inches(0.5),
         size=16, bold=True, color=GREEN)
add_text(s, "Only fear remains at F1 = 0.",
         Inches(8.3), Inches(2.0), Inches(4.8), Inches(0.5),
         size=16, bold=True, color=RED)

add_text(s, "Hypotheses",
         Inches(8.3), Inches(2.7), Inches(4.8), Inches(0.4),
         size=14, bold=True, color=NAVY)
bullets(s, [
    "Data-bound: 414 train utterances may be below the discrimination threshold.",
    "Appraisal overlap: fear signature overlaps with sadness + anger in the 8-d space.",
    "Capacity-bound: LoRA r = 8 may lack fine discriminative power for very rare classes.",
], Inches(8.3), Inches(3.15), Inches(4.8), Inches(3.5), size=12)
add_footer(s)

# ---- Slide 20: Ablations ------------------------------------------------
s = add_slide(); add_header(s, "Ablations — Surprising Findings", 20)

# Finding A
add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s, "(A)  Memory architecture matters more than presence of memory",
         Inches(0.7), Inches(1.45), Inches(12), Inches(0.5),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Stateless (no memory):   wF1 = 0.5631",
    "BiLSTM memory:                wF1 = 0.4241    ← 14 points WORSE than no memory.",
    "EmoFlow exp-decay:          wF1 = 0.6171    ← 19 points BETTER than BiLSTM.",
    "Generic recurrent priors actively hurt on small ERC datasets.",
], Inches(0.7), Inches(1.95), Inches(12), Inches(2.0), size=14)

# Finding B
add_rect(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(2.5), RGBColor(0xFF, 0xFF, 0xFF), line=MUTED)
add_text(s, "(B)  Learnable λ helps wF1 but slightly hurts rare-class metrics",
         Inches(0.7), Inches(4.2), Inches(12), Inches(0.5),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "λ = 0 (uniform):                    wF1 = 0.6052,   wF1_6way = 0.5348,   ETA = 0.5617",
    "learned λ:                              wF1 = 0.6171,   wF1_6way = 0.5312,   ETA = 0.5528",
    "Learned decay optimizes for the dominant class — but rare-class evidence may live in earlier turns.",
    "Implication: uniform aggregation may be preferable when rare classes are the priority.",
], Inches(0.7), Inches(4.7), Inches(12), Inches(2.0), size=14)
add_footer(s)

# ---- Slide 21: Interpretability case study -----------------------------
s = add_slide(); add_header(s, "Interpretability — Appraisal Case Study", 21)
add_text(s,
         "Test sentences not in training. We read the top-/bottom-activated appraisal dimensions.",
         Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5),
         size=15, color=INK)

data = [
    ["Sentence",                                          "Top dim",                       "Bottom dim",         "Match"],
    ["\"I am so happy! This is the best day ever.\"",     "self_consistency, coping",      "unpleasantness",     "Joy ✓"],
    ["\"Oh my god, what just happened?!\"",               "goal_hindrance, unpleasant.",   "expectedness",       "Surprise ✓"],
    ["\"Get the fuck out of my house!\"",                 "unfairness, immorality",        "self_consistency",   "Anger ✓"],
    ["\"I hate this. I want to die.\"",                   "unpleasant., goal_hindrance",   "self_consistency",   "Sad/Anger ✓"],
    ["\"Just another boring meeting today.\"",            "unpleasantness",                "coping_potential",   "Disgust/bored ✓"],
]
add_table(s, data, Inches(0.5), Inches(1.85), Inches(12.3), Inches(3.6),
          col_widths=[4.8, 3.5, 2.5, 1.5], font_size=12)

add_text(s, "Takeaway",
         Inches(0.5), Inches(5.65), Inches(12), Inches(0.4),
         size=16, bold=True, color=NAVY)
bullets(s, [
    "Encoder learned Scherer-aligned appraisal patterns WITHOUT per-sentence supervision.",
    "Only 6 emotion-level prototype vectors served as MSE targets.",
    "The 8-d bottleneck is genuinely interpretable — not a post-hoc rationalization.",
], Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.2), size=14)
add_footer(s)


# ===========================================================================
# PART IV — Wrap-up
# ===========================================================================
section_slide("IV", "What's Left & Who Did What", "Limitations, future work, and the team.")

# ---- Slide 22: Limitations ---------------------------------------------
s = add_slide(); add_header(s, "Limitations — What We Did Not Cover", 22)
bullets(s, [
    "(a)  Single primary dataset.  All main numbers on MELD. EmoryNLP results on the post-fix pipeline missing.  IEMOCAP not attempted.",
    "(b)  Response Generator not delivered.  Originally module 6 of the pipeline. Debugging budget consumed by §III.",
    "(c)  Human Coherence Score evaluation not performed.  Tied to (b).",
    "(d)  Fear class not learned.  F1 = 0 despite oversampling and rare-class augmentation.",
    "(e)  Bounded by frozen backbone.  ~20 M trainable can approach but not exceed fine-tuned 100 M+ models (COSMIC, MMGCN).",
    "(f)  Text only.  MELD has audio + video; prosody / facial expression cues are unused.",
    "(g)  Approximate Bayesian update.  Log-additive MLP fusion ≡ exact Dirichlet–Categorical only under uniform marginals.",
], Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7), size=15, gap=0.18)
add_footer(s)

# ---- Slide 23: Future work ---------------------------------------------
s = add_slide(); add_header(s, "Future Work — Given More Time", 23)
items = [
    ("1.  Complete the proposed pipeline",
     "Implement the Response Generator: emotion distribution → soft-prompt prefix → LLaMA-3-8B decoder.\n"
     "Evaluate via the Coherence Score protocol from the proposal."),
    ("2.  Multi-dataset evaluation",
     "Add EmoryNLP + IEMOCAP using the post-fix pipeline. The sigmoid-saturation diagnostic should transfer."),
    ("3.  Larger LoRA / partial unfreezing",
     "Test r = 16, 32 or unfreeze top LLaMA layers — target the fear-class gap specifically."),
    ("4.  Multimodal extension",
     "Inject MELD's audio prosodic features into the appraisal head — should help fear / sadness."),
    ("5.  Fully Bayesian inference",
     "Replace the log-additive MLP with an explicit Dirichlet–Categorical update, possibly variational."),
]
for i, (h, body) in enumerate(items):
    y = Inches(1.3 + i * 1.1)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.9), ACCENT)
    add_text(s, h, Inches(0.8), y, Inches(12), Inches(0.4),
             size=15, bold=True, color=NAVY)
    add_text(s, body, Inches(0.8), y + Inches(0.4), Inches(12), Inches(0.6),
             size=13, color=INK)
add_footer(s)

# ---- Slide 24: Team contributions --------------------------------------
s = add_slide(); add_header(s, "Team", 24)
add_text(s, "EmoFlow was built by",
         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
         size=18, color=MUTED, align=PP_ALIGN.CENTER)

names = ["Xiaoyi Liu", "Zhiye Jiang", "Ruitian Liu"]
for i, n in enumerate(names):
    x = Inches(1.0 + i * 4.0)
    add_rect(s, x, Inches(2.5), Inches(3.5), Inches(2.5), NAVY)
    add_text(s, n, x, Inches(3.1), Inches(3.5), Inches(0.8),
             size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, "UC Davis", x, Inches(3.9), Inches(3.5), Inches(0.5),
             size=14, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s,
         "All authors contributed jointly to the architecture, training, evaluation, and report.",
         Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
         size=16, color=INK, align=PP_ALIGN.CENTER)
add_text(s,
         "ECS 271 — Machine Learning — UC Davis, Spring 2026",
         Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
         size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s)


# ---- Closing thank-you (not counted) -----------------------------------
s = add_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
add_text(s, "Thank you",
         Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.2),
         size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, "Questions?",
         Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8),
         size=32, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "github.com / EmoFlow   ·   ECS 271 Spring 2026",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
         size=16, color=RGBColor(0xC0, 0xCB, 0xDC), align=PP_ALIGN.CENTER)


prs.save("EmoFlow_Presentation.pptx")
print("Saved EmoFlow_Presentation.pptx")
print(f"Total slides: {len(prs.slides)} (title + 4 section dividers + 24 content + closing)")